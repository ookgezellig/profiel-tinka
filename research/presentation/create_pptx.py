"""
Generate PPTX presentation for Tinka den Arend profile
Using PPF APG/APG styling (teal/petrol color scheme)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.dml.color import RGBColor
import os

# Color scheme based on PPF APG/APG branding
COLORS = {
    'primary': RGBColor(0x0D, 0x52, 0x57),      # Dark teal #0D5257
    'secondary': RGBColor(0x1a, 0x7a, 0x7f),    # Medium teal
    'accent': RGBColor(0xe8, 0xf4, 0xf5),       # Light teal background
    'text_dark': RGBColor(0x1F, 0x29, 0x37),    # Dark text
    'text_light': RGBColor(0xFF, 0xFF, 0xFF),   # White text
    'gray': RGBColor(0x6B, 0x72, 0x80),         # Gray for subtitles
}

# Paths
SCREENSHOTS_DIR = r"D:\KB-OPEN\github-repos\SaveToWaybackMachine\research\screenshots"
OUTPUT_DIR = r"D:\KB-OPEN\github-repos\SaveToWaybackMachine\research\presentation\output"

def set_slide_background(slide, color):
    """Set solid background color for a slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title, subtitle):
    """Add a title slide with dark background"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['primary'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_light']
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.33), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['text_light']
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, content_lines, image_path=None, image_right=True):
    """Add a content slide with optional image"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']

    # Determine content area based on image
    if image_path and os.path.exists(image_path):
        if image_right:
            content_left = Inches(0.5)
            content_width = Inches(6.5)
            img_left = Inches(7.5)
        else:
            content_left = Inches(6)
            content_width = Inches(6.5)
            img_left = Inches(0.5)

        # Add image
        try:
            slide.shapes.add_picture(image_path, img_left, Inches(1.5), width=Inches(5.5))
        except Exception as e:
            print(f"Could not add image {image_path}: {e}")
    else:
        content_left = Inches(0.5)
        content_width = Inches(12.33)

    # Content
    content_box = slide.shapes.add_textbox(content_left, Inches(1.5), content_width, Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Check if it's a bullet point
        if line.startswith('• '):
            p.text = line[2:]
            p.level = 0
        elif line.startswith('  - '):
            p.text = line[4:]
            p.level = 1
        else:
            p.text = line
            p.level = 0

        p.font.size = Pt(20)
        p.font.color.rgb = COLORS['text_dark']
        p.space_after = Pt(12)

    return slide

def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']

    # Table
    num_cols = len(headers)
    num_rows = len(rows) + 1  # +1 for header

    table_width = Inches(12)
    table_height = Inches(0.5 * num_rows)
    left = Inches(0.67)
    top = Inches(1.8)

    table = slide.shapes.add_table(num_rows, num_cols, left, top, table_width, table_height).table

    # Set column widths
    col_width = table_width / num_cols
    for i in range(num_cols):
        table.columns[i].width = int(col_width)

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['primary']
        p = cell.text_frame.paragraphs[0]
        p.font.color.rgb = COLORS['text_light']
        p.font.bold = True
        p.font.size = Pt(16)

    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS['text_dark']

    return slide

def add_quote_slide(prs, quote, attribution):
    """Add a slide with a quote"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['accent'])

    # Quote
    quote_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(3))
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f'"{quote}"'
    p.font.size = Pt(28)
    p.font.italic = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER

    # Attribution
    attr_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(11.33), Inches(0.5))
    tf = attr_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"— {attribution}"
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['gray']
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_section_slide(prs, title, subtitle=""):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['primary'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.33), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_light']
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(12.33), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['text_light']
        p.alignment = PP_ALIGN.CENTER

    return slide

def create_presentation():
    """Create the full presentation"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs,
        "Tinka den Arend",
        "Profiel van een Pensioenpionier\n\nWerkgeversvoorzitter PPF APG | Hoofd Policy & Public Affairs APG")

    # Slide 2: Profile Summary
    profile_photo = os.path.join(SCREENSHOTS_DIR, "tinka_profile_photo.jpg")
    add_content_slide(prs, "Profielsamenvatting", [
        "Tinka den Arend is een sleutelfiguur in de Nederlandse pensioentransitie",
        "",
        "• Hoofd team Policy and Public Affairs bij APG",
        "• Werkgeversvoorzitter van PPF APG sinds 2019",
        "• Brede achtergrond in publieke sector",
        "• Leidde PPF APG als 'koploper' naar nieuw pensioenstelsel",
    ], profile_photo)

    # Slide 3: Current Positions (table)
    add_table_slide(prs, "Huidige Functies",
        ["Functie", "Organisatie", "Sinds"],
        [
            ["Hoofd team Policy and Public Affairs", "APG", "Huidig"],
            ["Bestuurslid & Werkgeversvoorzitter", "PPF APG", "2019"],
        ])

    # Slide 4: Career Overview
    add_content_slide(prs, "Carrière Overzicht", [
        "2019 - heden: Bestuurslid & Werkgeversvoorzitter PPF APG",
        "2016 - 2019: Lid/voorzitter Verantwoordingsorgaan PPF APG",
        "Eerder: Diverse functies bij Tweede Kamer",
        "Eerder: Consultant/Adviseur bij PwC",
        "Eerder: Beleidsmedewerker Min. van SZW",
        "Eerder: Medewerker LISV",
        "",
        "Rode draad: Pensioenwetgeving, politieke processen, uitvoering"
    ])

    # Slide 5: Expertise
    add_table_slide(prs, "Expertisegebieden",
        ["Domein", "Beschrijving"],
        [
            ["Pensioenwetgeving", "Wet toekomst pensioenen (Wtp)"],
            ["Pensioentransitie", "Invaren, solidariteitsreserve"],
            ["Governance", "Bestuursstructuren pensioenfondsen"],
            ["Beleidsadvies", "Strategische advisering RvB"],
            ["Sociale zekerheid", "Achtergrond bij SZW en LISV"],
            ["Politieke processen", "Ervaring Tweede Kamer"],
        ])

    # Slide 6: Section - 2025
    add_section_slide(prs, "2025: Het Jaar van de Doorbraak",
        "PPF APG als Koploper in de Pensioentransitie")

    # Slide 7: Milestones 2025
    add_content_slide(prs, "Belangrijkste Mijlpalen 2025", [
        "Februari 2025:",
        "• Eerste pensioenen uitbetaald onder nieuw stelsel",
        "• PPF APG was een van de drie 'koplopers'",
        "",
        "25 juni 2025:",
        "• Pensioen Pro Award 'Pensioenfonds van het Jaar'",
        "",
        "De Drie Koplopers:",
        "• PPF APG, PWRI, Beroepspensioenfonds Loodsen",
        "• 65.000 gepensioneerden ontvingen eerste uitkering"
    ])

    # Slide 8: Award with screenshot
    award_img = os.path.join(SCREENSHOTS_DIR, "10_ppf_apg_nieuws_award.png")
    add_content_slide(prs, "Award: Pensioenfonds van het Jaar", [
        "De jury prees PPF APG's voorbeeldrol als koploper",
        "",
        "Citaat Tinka den Arend:",
        "",
        "'Ik ben trots op de samenwerking met onze",
        "sociale partners, het verantwoordingsorgaan,",
        "de Raad van Toezicht, het bestuur en de",
        "APG-medewerkers die de transitie mogelijk",
        "hebben gemaakt.'",
        "",
        "— juni 2025"
    ], award_img)

    # Slide 9: Speaking engagements
    add_table_slide(prs, "Spreekbeurten & Media 2025",
        ["Datum", "Evenement", "Details"],
        [
            ["6 feb", "Instituut Pensioen Educatie", "Presentatie transitie-ervaringen"],
            ["14 feb", "Pensioen Pro Podcast", "Interview 'Twee keer invaren' (44 min)"],
            ["2025", "PensioenPro Focuscongres", "Spreker pensioentransitie"],
            ["6 nov", "CFA Society ALM Congres", "Panellid beleggingen nieuw stelsel"],
            ["Okt", "Gepensioneerdenbijeenkomst", "Presentatie aan deelnemers"],
        ])

    # Slide 10: Podcast with screenshot
    podcast_img = os.path.join(SCREENSHOTS_DIR, "05_pensioenpro_podcast.png")
    add_content_slide(prs, "Podcast: 'Twee Keer Invaren'", [
        "Pensioen Pro In Gesprek",
        "14 februari 2025 | 44 minuten",
        "",
        "Besproken onderwerpen:",
        "• Het unieke 'twee keer invaren' proces",
        "• Haar carrièrepad naar pensioensector",
        "• Politieke actualiteit via APG",
        "• Beleggingsstrategie nieuw stelsel",
        "",
        "'Heel erg soepel gegaan'"
    ], podcast_img)

    # Slide 11: CFA Congress with screenshot
    cfa_img = os.path.join(SCREENSHOTS_DIR, "07_cfa_alm_congres.png")
    add_content_slide(prs, "CFA Society ALM Congres 2025", [
        "6 november 2025",
        "",
        "Thema: ALM in Transitie",
        "• Invaren",
        "• Innovatie",
        "• Internationaal",
        "",
        "Tinka den Arend treedt op als panellid",
        "in de sessie over 'Beleggen in het nieuwe",
        "pensioenstelsel'"
    ], cfa_img)

    # Slide 12: What is "Twee Keer Invaren"
    add_content_slide(prs, "Wat is 'Twee Keer Invaren'?", [
        "PPF APG had een unieke situatie:",
        "het fonds moest twee keer invaren vanwege",
        "zijn specifieke regelingsstructuur.",
        "",
        "Proces:",
        "1. Eerste invaring: Bestaande rechten omzetten",
        "2. Tweede invaring: Aanvullende regelingscomponenten",
        "",
        "Dit maakte PPF APG tot een waardevolle",
        "testcase voor de rest van de pensioensector."
    ])

    # Slide 13: Board tasks
    add_table_slide(prs, "Bestuurlijke Taken bij PPF APG",
        ["Commissie/Rol", "Verantwoordelijkheid"],
        [
            ["Commissie balansbeheer en beleggingen", "Toezicht op beleggingsbeleid"],
            ["Verkiezingscommissie", "Organisatie bestuursverkiezingen"],
            ["Strategie van het fonds", "Langetermijnplanning"],
        ])

    # Slide 14: PPF APG Profile with screenshot
    profile_img = os.path.join(SCREENSHOTS_DIR, "01_ppf_apg_tinka_profiel.png")
    add_content_slide(prs, "PPF APG Profiel", [
        "Officiële profielpagina op ppf-apg.nl",
        "",
        "Het Personeelspensioenfonds APG is het",
        "pensioenfonds voor medewerkers van APG zelf.",
        "",
        "DNB Registratie:",
        "• Statutaire naam: Stichting Personeelspensioenfonds APG",
        "• DNB Relatienummer: 02392"
    ], profile_img)

    # Slide 15: Publications
    add_table_slide(prs, "Publicaties onder haar Voorzitterschap",
        ["Document", "Jaar"],
        [
            ["PPF APG Jaarverslag 2024", "juni 2025"],
            ["Verslag Verantwoord Beleggen 2024", "2025"],
            ["Aandelenportefeuille PPF APG", "juni 2025"],
            ["PPF APG Jaarverslag 2023", "juni 2024"],
            ["PPF APG Jaarverslag 2022", "2023"],
            ["PPF APG Jaarverslag 2021", "2022"],
        ])

    # Slide 16: Board with screenshot
    bestuur_img = os.path.join(SCREENSHOTS_DIR, "08_ppf_apg_bestuur.png")
    add_content_slide(prs, "Bestuur PPF APG", [
        "Het bestuur van PPF APG bestaat uit",
        "werkgevers- en werknemersvertegenwoordigers.",
        "",
        "Tinka den Arend fungeert als",
        "werkgeversvoorzitter en geeft leiding",
        "aan de strategische richting van het fonds."
    ], bestuur_img)

    # Slide 17: Online presence
    add_table_slide(prs, "Online Aanwezigheid",
        ["Platform", "Link"],
        [
            ["LinkedIn", "nl.linkedin.com/in/tinka-den-arend-3b218b11"],
            ["PPF APG Profiel", "ppf-apg.nl/bestuur/tinka-den-arend"],
            ["The Org", "theorg.com/org/werkenbijapg"],
            ["DNB Register", "dnb.nl/public-register"],
        ])

    # Slide 18: Section - Summary
    add_section_slide(prs, "Samenvatting")

    # Slide 19: Key Points
    add_content_slide(prs, "Kernpunten", [
        "Wie is Tinka den Arend?",
        "• Ervaren professional met brede publieke sector achtergrond",
        "• Dubbele functie: APG (Policy & Public Affairs) + PPF APG",
        "• Expert in pensioenwetgeving, governance en transitie",
        "",
        "Prestaties 2025:",
        "• Koploper in pensioentransitie naar nieuw stelsel",
        "• Pensioen Pro Award voor Pensioenfonds van het Jaar",
        "• Thought leader via podcasts, congressen en publicaties"
    ])

    # Slide 20: Sources
    add_content_slide(prs, "Bronnen", [
        "• PPF APG website (ppf-apg.nl)",
        "• APG website (apg.nl)",
        "• Pensioen Pro (pensioenpro.nl)",
        "• CFA Society Netherlands (cfasociety.nl)",
        "• DNB Register (dnb.nl)",
        "• LinkedIn, The Org",
        "",
        "Onderzoeksdatum: 1 december 2025",
        "Onderzoeksmethode: Exa MCP deep search"
    ])

    # Slide 21: Questions
    add_title_slide(prs, "Vragen?",
        "Contact: Via officiële kanalen APG en PPF APG")

    # Save
    output_path = os.path.join(OUTPUT_DIR, "tinka_den_arend_profiel.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
